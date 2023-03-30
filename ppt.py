import collections 
import collections.abc
import json
from pptx import Presentation

def ppt(username):
    prs = Presentation()

    # Read the JSON object from a file
    with open(username+"ppt1.json", 'r') as f:
        ppt = json.load(f)

    # Extract the title and slides fields from the JSON object
    title = ppt['title']
    slides = ppt['slides']

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title1 = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title1.text = title
    subtitle.text = "By ChatGPT"

    # Loop through each slide and extract the title and content fields
    for slide in slides:
        slide_title = slide['title']
        slide_content = slide['content']
        print(f"Slide title: {slide_title}")
        print("Slide content:")

        # Add a bullet slide
        bullet_slide_layout = prs.slide_layouts[3]
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        title2 = slide2.shapes.title
        body2 = slide2.shapes.placeholders[1]
        title2.text = slide_title
        tf = body2.text_frame

        for content_string in slide_content:
            p = tf.add_paragraph()
            p.text = content_string
            p.level = 1

        print("\n\n")



    prs.save(username+str('.pptx'))